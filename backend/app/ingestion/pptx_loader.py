# app/ingestion/pptx_loader.py
from __future__ import annotations

from pathlib import Path
from typing import List, Dict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _clean(s: str) -> str:
    return (s or "").strip()


def load_pptx(path: Path, course_id: str, doc_name: str) -> List[Dict]:
    """
    Extract per slide:
      - text boxes/titles
      - tables
      - image alt-text/name (diagram hints)
      - speaker notes (often the most useful)
    """
    prs = Presentation(str(path))
    chunks: List[Dict] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []

        # Slide shapes
        for shape in slide.shapes:
            # Text frames
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                t = _clean(shape.text_frame.text)
                if t:
                    texts.append(t)

            # Tables
            if getattr(shape, "has_table", False) and shape.has_table:
                tbl = shape.table
                rows = []
                for r in tbl.rows:
                    row_cells = [_clean(c.text) for c in r.cells]
                    if any(row_cells):
                        rows.append(" | ".join(row_cells))
                if rows:
                    texts.append("[Table]\n" + "\n".join(rows))

            # Pictures (diagrams)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                alt = _clean(getattr(shape, "alternative_text", "") or "")
                name = _clean(getattr(shape, "name", "") or "")
                hint = alt or name
                if hint:
                    texts.append(f"[Diagram/Image hint]\n{hint}")

        # Speaker notes (very useful)
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes = _clean(slide.notes_slide.notes_text_frame.text)
                if notes:
                    texts.append(f"[Speaker Notes]\n{notes}")
        except Exception:
            pass

        if not texts:
            continue

        full_text = "\n\n".join(texts)
        chunks.append(
            {
                "course_id": course_id,
                "doc_name": doc_name,
                "source_type": "pptx",
                "page_or_slide": slide_idx,
                "text": full_text,
            }
        )

    return chunks